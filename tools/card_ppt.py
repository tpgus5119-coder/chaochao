#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""카드뉴스를 **고칠 수 있는 파워포인트**로 내보낸다 → 바탕화면 날짜 폴더

대표님 지시 (2026-08-31): "니가 만든 카드 내가 수정할수잇도록 ppt형식으로 해줄수도 잇니?"

webp 카드는 글자가 그림에 구워져 있어 못 고친다. 그래서 파워포인트는 이렇게 만든다:
  · **막까지 입힌 배경**(img/card/bg/*.png — card_news.py 가 글자 얹기 직전에 남긴다)을 바닥에 깐다
  · 갈래·제목·본문·낱말을 전부 **글상자**로 얹는다 → 파워포인트에서 그대로 고친다
파워포인트에서 반투명 도형을 따로 얹지 않는다 — 배경에 이미 막이 입혀져 있어
결과가 webp 카드와 똑같고, 도형 하나가 줄어 고치기도 쉽다.
슬라이드는 카드와 같은 정사각형(1080×1080 → 11.25인치)이다.

쓰기: python3 tools/card_ppt.py [--day 2026-08-28]
"""
import argparse, json, pathlib, sys

R = pathlib.Path(__file__).resolve().parent.parent
sys.path.insert(0, str(R / "tools"))
import vi_kr
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

CARD = R / "img" / "card"
DESK = pathlib.Path.home() / "Desktop" / "chaochao-cardnews"

PX = 11.25 / 1080                      # 카드 1080px = 11.25인치
E = lambda px: Emu(int(px * PX * 914400))
# 카드의 픽셀 글꼴을 그대로 포인트로 옮긴다. 1px = 11.25인치/1080 = 0.75pt.
# (이 환산을 안 하면 파워포인트 글자가 카드보다 훨씬 작아진다 — 2026-08-31에 실제로 그랬다)
PT = lambda px: Pt(round(px * 0.75, 1))

BG   = RGBColor(0xFA, 0xF9, 0xF6)
FG   = RGBColor(0x18, 0x1A, 0x1E)
BODY = RGBColor(0x26, 0x2A, 0x32)
DIM  = RGBColor(0x6E, 0x73, 0x7D)
CAT_COLOR = {"일자리": (33, 90, 160), "공장·산업": (120, 60, 20), "경제": (20, 90, 70),
             "사회": (130, 40, 90), "정치": (60, 60, 60), "문화·생활": (150, 90, 20)}
KO, VI = "Apple SD Gothic Neo", "Helvetica"    # 맥 파워포인트·키노트가 바로 찾는 이름


def box(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(E(x), E(y), E(w), E(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def put(tf, text, size_px, color, bold=False, font=KO, space_after=0, first=False, lh=1.35):
    para = tf.paragraphs[0] if first else tf.add_paragraph()
    r = para.add_run(); r.text = text
    r.font.size = PT(size_px); r.font.bold = bold
    r.font.color.rgb = color; r.font.name = font
    para.space_after = PT(space_after)
    para.line_spacing = lh
    return para


def slide1(prs, d, ts, i):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid(); s.background.fill.fore_color.rgb = BG
    # 배경은 **막까지 입힌 상태**를 card_news.py 가 png 로 남겨 둔다.
    # 파워포인트에서 반투명 도형을 얹는 것보다 이 편이 결과가 카드와 똑같다.
    bg = CARD / "bg" / f"{ts}-{i}.png"
    if bg.exists():
        s.shapes.add_picture(str(bg), 0, 0, E(1080), E(1080))

    cat = d.get("cat") or "소식"
    from pptx.enum.shapes import MSO_SHAPE
    chip = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, E(62), E(58), E(200), E(56))
    chip.fill.solid(); chip.fill.fore_color.rgb = RGBColor(*CAT_COLOR.get(cat, (60, 60, 60)))
    chip.line.fill.background(); chip.shadow.inherit = False
    ctf = chip.text_frame; ctf.word_wrap = False
    ctf.paragraphs[0].alignment = PP_ALIGN.CENTER
    cr = ctf.paragraphs[0].add_run(); cr.text = cat
    cr.font.size = PT(30); cr.font.bold = True; cr.font.name = KO
    cr.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    tf = box(s, 62, 140, 956, 200)
    # **카드 그림과 똑같은 글을 쓴다** (대표님 지적 2026-09-02 "pptx 랑 webp 랑 내용이 다르잖아").
    put(tf, d.get("title_card") or d.get("title", ""), 50, FG, bold=True, first=True, lh=1.30)

    tf2 = box(s, 62, 330, 956, 640)
    lines = d.get("sum5") or [d.get("intro") or ""]
    for k, ln in enumerate(lines):
        put(tf2, ln, 30, BODY, space_after=22, first=(k == 0), lh=1.30)

    put(box(s, 62, 990, 700, 40), foot_text(d, ts), 24, DIM, first=True)
    return s


def foot_text(d, ts):
    """카드 그림과 **똑같은** 아래 줄 — 펴낸 날 · 출처"""
    import urllib.parse
    host = urllib.parse.urlparse(d.get("u") or "").netloc.lower().replace("www.", "")
    name = {"insidevina.com": "인사이드비나", "vnexpress.net": "VnExpress",
            "e.vnexpress.net": "VnExpress", "vietnamkoreatimes.com": "베트남코리아타임즈",
            "tuoitre.vn": "Tuổi Trẻ", "thanhnien.vn": "Thanh Niên"}.get(host, host)
    return "짜오짜오 · " + (d.get("pub") or ts) + ("  ·  출처 " + name if name else "")


def slide2(prs, d, ts):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid(); s.background.fill.fore_color.rgb = BG
    # 카드 그림과 같은 짜임 — 제목 없이 낱말 여섯, 아래에 대화 두 줄
    words = (d.get("words") or [])[:6]
    for i, w in enumerate(words):
        x = 84 + (i % 2) * 484
        y = 80 + (i // 2) * 190
        t = box(s, x, y, 452, 180)
        put(t, w["vi"], 44, FG, bold=True, font=VI, first=True)
        # 발음은 **늘 우리 도구**가 만든다 (AI 것은 안 쓴다)
        put(t, "[" + (vi_kr.word(w["vi"]) or w.get("kr_read") or "") + "]", 25, DIM, space_after=10)
        put(t, w.get("ko", ""), 29, BODY)
    lines = ((d.get("dialog") or {}).get("lines") or [])[:2]
    if lines:
        h = box(s, 84, 700, 912, 60)
        put(h, "이 기사로 나누는 말", 34, FG, bold=True, first=True)
        y = 770
        for ln in lines:
            t = box(s, 84, y, 912, 120)
            put(t, (ln.get("who") or "A") + ". " + (ln.get("vi") or ""), 31, FG,
                bold=True, font=VI, first=True)
            put(t, ln.get("ko") or "", 27, BODY, space_after=10)
            y += 100
    put(box(s, 84, 990, 700, 40), foot_text(d, ts), 24, DIM, first=True)
    return s


def main():
    a = argparse.ArgumentParser(); a.add_argument("--day"); a = a.parse_args()
    days = json.loads((R / "data" / "news_days.json").read_text(encoding="utf-8"))["days"]
    # **펴낸 날로 묶는다.** 기사 날짜로 묶으면 이틀치가 한 폴더에 올 때
    # 파워포인트가 하루치만 담긴다 (실측 2026-09-02: 22장 중 8장만 들어갔다).
    by = {}
    for d in days:
        by.setdefault(d.get("pub") or d.get("ts"), []).append(d)

    n = 0
    for pub, arts in sorted(by.items()):
        if not pub or (a.day and pub != a.day): continue
        # 그림 파일 번호는 **기사 날짜 안에서** 매겨진다
        cnt = {}
        for d in arts:
            k = d.get("ts"); cnt[k] = cnt.get(k, 0) + 1; d["_i"] = cnt[k]
        arts = [d for d in arts if (CARD / f"{d.get('ts')}-{d['_i']}-1.webp").exists()]
        if not arts: continue
        prs = Presentation()
        prs.slide_width, prs.slide_height = E(1080), E(1080)
        for d in arts:
            slide1(prs, d, d.get("ts"), d["_i"])
            slide2(prs, d, d.get("ts"))
        out = DESK / pub
        out.mkdir(parents=True, exist_ok=True)
        f = out / f"카드뉴스-{pub}.pptx"
        prs.save(str(f))
        print(f"  {pub}: 슬라이드 {len(arts) * 2}장 → {f}")
        n += 1
    print(f"파워포인트 {n}개 만들었습니다" if n else "만들 것이 없습니다")


main()
